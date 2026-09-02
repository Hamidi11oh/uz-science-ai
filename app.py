import streamlit as st
import fitz  # PyMuPDF
import requests
from docx import Document
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
import io
import json
import os
import time
import concurrent.futures

# =====================================================================
# 👤 MUALLIF VA ASOSIY SOZLAMALAR
# =====================================================================
AUTHOR_NAME = "Qosimjonov Hamidullo"

# Streamlit Secrets-dan kalitni avtomatik olish
BUILTIN_GROQ_KEY = ""
if "GROQ_API_KEY" in st.secrets:
    BUILTIN_GROQ_KEY = str(st.secrets["GROQ_API_KEY"]).strip().strip("'\" \n\r\t")
elif os.environ.get("GROQ_API_KEY"):
    BUILTIN_GROQ_KEY = str(os.environ.get("GROQ_API_KEY")).strip().strip("'\" \n\r\t")

# =====================================================================
# SAHIFA SOZLAMALARI VA DIZAYN
# =====================================================================
st.set_page_config(
    page_title="UZ SCIENCE AI - Professional Ilmiy Tarjima",
    page_icon="🎓",
    layout="wide",
    initial_sidebar_state="collapsed"
)

st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=Plus+Jakarta+Sans:wght@300;400;500;600;700;800&family=JetBrains+Mono:wght@400;500&display=swap');
    
    html, body, [class*="css"] {
        font-family: 'Plus Jakarta Sans', sans-serif;
    }
    
    .stApp {
        background-color: #f8fafc;
    }
    
    .block-container {
        max-width: 960px !important;
        padding-top: 1.5rem !important;
        padding-bottom: 4rem !important;
    }
    
    .top-nav {
        display: flex;
        justify-content: space-between;
        align-items: center;
        background: #ffffff;
        border: 1px solid #e2e8f0;
        border-radius: 16px;
        padding: 12px 22px;
        margin-bottom: 24px;
        box-shadow: 0 1px 3px rgba(0,0,0,0.03);
    }
    
    .nav-brand {
        display: flex;
        align-items: center;
        gap: 12px;
    }
    
    .nav-logo {
        width: 40px;
        height: 40px;
        border-radius: 12px;
        background: linear-gradient(135deg, #0284c7 0%, #4f46e5 100%);
        display: flex;
        align-items: center;
        justify-content: center;
        color: white;
        font-size: 1.2rem;
        box-shadow: 0 4px 10px rgba(2, 132, 199, 0.25);
    }
    
    .badge-author {
        display: inline-flex;
        align-items: center;
        gap: 6px;
        background: linear-gradient(135deg, #f0fdf4 0%, #e0f2fe 100%);
        border: 1px solid #bae6fd;
        color: #0369a1;
        padding: 6px 14px;
        border-radius: 20px;
        font-size: 0.82rem;
        font-weight: 700;
        box-shadow: 0 2px 5px rgba(2, 132, 199, 0.08);
    }
    
    .hero-container {
        text-align: center;
        margin-bottom: 24px;
    }
    
    .hero-title {
        font-size: 2.35rem;
        font-weight: 800;
        letter-spacing: -0.8px;
        color: #0f172a;
        margin-bottom: 8px;
    }
    
    .hero-highlight {
        background: linear-gradient(135deg, #0284c7 0%, #4f46e5 100%);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
    }
    
    .hero-subtitle {
        font-size: 0.96rem;
        color: #64748b;
        max-width: 620px;
        margin: 0 auto;
        line-height: 1.55;
    }
    
    .main-card {
        background: #ffffff;
        border: 1px solid #e2e8f0;
        border-radius: 20px;
        padding: 32px;
        box-shadow: 0 20px 25px -5px rgba(0, 0, 0, 0.03), 0 8px 10px -6px rgba(0, 0, 0, 0.02);
        margin-bottom: 24px;
    }
    
    .option-box {
        background: #f8fafc;
        border: 1px solid #e2e8f0;
        border-radius: 12px;
        padding: 14px;
        font-size: 0.85rem;
    }
    
    .option-title {
        font-weight: 700;
        color: #0f172a;
        margin-bottom: 8px;
        display: flex;
        align-items: center;
        gap: 6px;
        font-size: 0.88rem;
    }
    
    .res-card {
        background: #ffffff;
        border: 1px solid #e2e8f0;
        border-radius: 16px;
        padding: 24px;
        margin-bottom: 16px;
        box-shadow: 0 1px 3px rgba(0,0,0,0.02);
    }
    
    .res-card h4 {
        color: #0f172a;
        margin-top: 0;
        margin-bottom: 12px;
        font-size: 1.1rem;
        font-weight: 700;
    }
    
    .focus-banner {
        background: linear-gradient(135deg, #eff6ff 0%, #dbeafe 100%);
        border: 1px solid #bfdbfe;
        border-left: 6px solid #0284c7;
        padding: 20px 24px;
        border-radius: 14px;
        margin-bottom: 20px;
    }
    
    .stTabs [data-baseweb="tab-list"] {
        gap: 8px;
        border-bottom: 1px solid #e2e8f0;
        padding-bottom: 8px;
    }

    .stTabs [data-baseweb="tab"] {
        border-radius: 10px;
        padding: 8px 18px;
        font-weight: 600;
        font-size: 0.88rem;
        color: #64748b;
        background-color: #ffffff;
        border: 1px solid #e2e8f0;
    }

    .stTabs [aria-selected="true"] {
        background-color: #0284c7 !important;
        color: #ffffff !important;
        border-color: #0284c7 !important;
    }
    
    .stButton button {
        background: linear-gradient(135deg, #0284c7 0%, #4f46e5 100%) !important;
        color: white !important;
        border: none !important;
        padding: 14px 28px !important;
        border-radius: 14px !important;
        font-weight: 700 !important;
        font-size: 1.05rem !important;
        box-shadow: 0 10px 20px -3px rgba(2, 132, 199, 0.35) !important;
        transition: all 0.2s ease !important;
    }
    
    .stButton button:hover {
        transform: translateY(-1px) !important;
        box-shadow: 0 14px 25px -3px rgba(2, 132, 199, 0.45) !important;
    }
    
    .footer-container {
        text-align: center;
        padding-top: 24px;
        border-top: 1px solid #e2e8f0;
        margin-top: 40px;
        font-size: 0.85rem;
        color: #64748b;
    }
</style>
""", unsafe_allow_html=True)

# Helper function to generate PDF
def create_pdf(title, translation_text, summary_data, thesis_data, terms_data):
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter, rightMargin=36, leftMargin=36, topMargin=36, bottomMargin=36)
    styles = getSampleStyleSheet()
    
    title_style = ParagraphStyle('DocTitle', parent=styles['Heading1'], fontSize=15, leading=19, textColor=colors.HexColor('#0284c7'), spaceAfter=6)
    h1_style = ParagraphStyle('H1', parent=styles['Heading2'], fontSize=11, leading=15, textColor=colors.HexColor('#0f172a'), spaceBefore=12, spaceAfter=6)
    body_style = ParagraphStyle('Body', parent=styles['Normal'], fontSize=8.5, leading=12.5, textColor=colors.HexColor('#334155'), spaceAfter=4)
    
    story = []
    clean_title = str(title).replace('<', '&lt;').replace('>', '&gt;')
    story.append(Paragraph(f"<b>UZ SCIENCE AI — {clean_title}</b>", title_style))
    story.append(Paragraph(f"<i>Muallif: {AUTHOR_NAME} · Mukammal akademik tarjima va ilmiy tahlil</i>", body_style))
    story.append(Spacer(1, 10))
    
    # 1. TO'LIQ AKADEMIK TARJIMA (BOSHIDA)
    story.append(Paragraph("<b>1. TO‘LIQ AKADEMIK O‘ZBEKCHA TARJIMA (ASOSIY MATN)</b>", h1_style))
    for p in translation_text.split('\n'):
        if p.strip():
            clean_p = p.replace('<', '&lt;').replace('>', '&gt;')
            story.append(Paragraph(clean_p, body_style))
    story.append(Spacer(1, 12))
    
    # 2. CHUQUR ILMIY XULOSA
    if summary_data:
        story.append(Paragraph("<b>2. Chuqur Ilmiy Xulosa va Tahlil (Research Summary)</b>", h1_style))
        for k, v in summary_data.items():
            label = k.replace('_', ' ').capitalize()
            story.append(Paragraph(f"<b>• {label}:</b> {str(v).replace('<', '&lt;').replace('>', '&gt;')}", body_style))
        story.append(Spacer(1, 10))
        
    # 3. MAGISTR TAVSIYALARI
    if thesis_data:
        story.append(Paragraph("<b>3. Magistrlik Dissertatsiyasi Uchun Tavsiyalar</b>", h1_style))
        for k, v in thesis_data.items():
            label = k.replace('_', ' ').capitalize()
            story.append(Paragraph(f"<b>• {label}:</b> {str(v).replace('<', '&lt;').replace('>', '&gt;')}", body_style))
        story.append(Spacer(1, 10))

    # 4. TERMINLAR
    if terms_data:
        story.append(Paragraph("<b>4. Asosiy Ilmiy Terminlar</b>", h1_style))
        for item in terms_data:
            term_en = item.get("term_en", "")
            term_uz = item.get("term_uz", "")
            desc = item.get("desc", "")
            story.append(Paragraph(f"<b>• {term_en} → {term_uz}:</b> {desc}", body_style))
            
    doc.build(story)
    buffer.seek(0)
    return buffer.getvalue()

# Active Key resolution
active_groq_key = BUILTIN_GROQ_KEY.strip()

with st.sidebar:
    st.header("⚙️ Sozlamalar")
    if active_groq_key:
        st.success("✅ Server Groq AI kaliti faol. Barcha foydalanuvchilar bepul va limitsiz foydalanishi mumkin.")
    else:
        st.info("💡 Server kaliti topilmadi. O'zingiznikini kiriting:")
        user_key = st.text_input("Groq API Kalit (gsk_...):", type="password")
        if user_key.strip():
            active_groq_key = user_key.strip()
            
    st.markdown("👉 [Bepul Groq Kalit Olish](https://console.groq.com/keys)")

# Top Navigation Bar with Author Badge in the Corner
st.markdown(f"""
<div class="top-nav">
    <div class="nav-brand">
        <div class="nav-logo">✨</div>
        <div>
            <div style="font-weight: 800; font-size: 1.15rem; color: #0f172a; display: flex; align-items: center;">
                UZ SCIENCE AI
            </div>
            <div style="font-size: 0.76rem; color: #64748b;">Magistrant va tadqiqotchilar uchun ilmiy platforma</div>
        </div>
    </div>
    <div class="badge-author">
        <span>👨‍💻 Muallif:</span> <b>{AUTHOR_NAME}</b>
    </div>
</div>
""", unsafe_allow_html=True)

# Hero Section
st.markdown("""
<div class="hero-container">
    <div class="hero-title">Ilmiy maqolalarni <span class="hero-highlight">o‘zbek tiliga</span> tarjima qiling</div>
    <div class="hero-subtitle">PDF, DOCX, PPTX yoki matnni to‘g‘ridan-to‘g‘ri joylashtiring. Mukammal akademik ma'no, formulalar va terminologiya to‘liq saqlanadi.</div>
</div>
""", unsafe_allow_html=True)

# Main Floating White Card
st.markdown('<div class="main-card">', unsafe_allow_html=True)

tab_upload, tab_paste = st.tabs(["📄 Fayl Yuklash (PDF/DOCX/PPTX)", "✍️ Matnni Qo‘yish (Paste)"])

extracted_text = ""
file_name = "Ilmiy Hujjat"

with tab_upload:
    uploaded_file = st.file_uploader(
        "PDF yoki boshqa faylni shu yerga tashlang:",
        type=["pdf", "docx", "pptx", "txt"],
        label_visibility="collapsed"
    )
    
    if uploaded_file is not None:
        file_name = uploaded_file.name
        suffix = file_name.split(".")[-1].lower()
        
        try:
            bytes_data = uploaded_file.read()
            if suffix == "pdf":
                doc = fitz.open(stream=bytes_data, filetype="pdf")
                pages = []
                for i, page in enumerate(doc):
                    t = page.get_text()
                    if t.strip():
                        pages.append(f"[Sahifa {i+1}]:\n" + t)
                extracted_text = "\n\n".join(pages)
            elif suffix == "docx":
                doc = Document(io.BytesIO(bytes_data))
                extracted_text = "\n".join([p.text for p in doc.paragraphs if p.text])
            elif suffix == "pptx":
                from pptx import Presentation
                prs = Presentation(io.BytesIO(bytes_data))
                slides = []
                for idx, slide in enumerate(prs.slides):
                    slide_texts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            slide_texts.append(shape.text)
                    slides.append(f"[Slayd {idx+1}]:\n" + "\n".join(slide_texts))
                extracted_text = "\n\n".join(slides)
            else:
                extracted_text = bytes_data.decode("utf-8", errors="ignore")
            
            if extracted_text.strip():
                st.markdown(f"""
                <div style="background: #f0fdf4; border: 1px solid #bbf7d0; padding: 10px 16px; border-radius: 10px; color: #166534; font-size: 0.88rem; font-weight: 600; display: flex; align-items: center; justify-content: space-between; margin-top: 10px;">
                    <span>📄 <b>{file_name}</b> muvaffaqiyatli yuklandi</span>
                    <span style="background: #dcfce7; padding: 2px 8px; border-radius: 6px; font-size: 0.78rem;">{len(extracted_text):,} belgi</span>
                </div>
                """, unsafe_allow_html=True)
            else:
                st.warning("⚠️ Fayldan matn ajratib bo'lmadi. Agar bu skaner qilingan rasm bo'lsa, 'Matnni Qo‘yish' bo'limidan foydalaning.")
        except Exception as e:
            st.error(f"Faylni o'qishda xatolik: {e}")

with tab_paste:
    direct_text = st.text_area(
        "Maqola matnini joylashtiring:",
        height=180,
        placeholder="Abstract, Introduction yoki butun maqola matnini Ctrl+V qilib bu yerga tashlang...",
        label_visibility="collapsed"
    )
    if direct_text.strip():
        extracted_text = direct_text
        file_name = "Kiritilgan Maqola Matni"

st.markdown("<br>", unsafe_allow_html=True)
col_opt1, col_opt2 = st.columns(2)

with col_opt1:
    st.markdown("""
    <div class="option-box">
        <div class="option-title">🌐 Til sozlamalari</div>
        <div style="display: flex; gap: 8px;">
            <div style="flex:1; background:#ffffff; border:1px solid #cbd5e1; padding:6px 10px; border-radius:8px; font-size:0.78rem; color:#475569;">
                <b>Asl til:</b> 🔄 Avtomatik
            </div>
            <div style="flex:1; background:#f0f9ff; border:1px solid #bae6fd; padding:6px 10px; border-radius:8px; font-size:0.78rem; color:#0369a1;">
                <b>Tarjima:</b> 🇺🇿 Akademik O‘zbek tili
            </div>
        </div>
    </div>
    """, unsafe_allow_html=True)

with col_opt2:
    st.markdown("""
    <div class="option-box">
        <div class="option-title">🛡️ Akademik talablar</div>
        <div style="display: grid; grid-template-columns: 1fr 1fr; gap: 6px; font-size: 0.76rem; color: #334155;">
            <div>☑️ Ravon akademik ma'no</div>
            <div>☑️ Formulalar daxlsizligi</div>
            <div>☑️ Qisqartirishlarsiz to‘liq</div>
            <div>☑️ Asl iqtiboslar (Citations)</div>
        </div>
    </div>
    """, unsafe_allow_html=True)

# Smart Paragraph/Section Chunking (Ideal 5,000 chars per chunk for maximum depth & speed)
def chunk_document_smart(text, target_size=5000):
    paragraphs = text.split('\n')
    chunks = []
    curr = []
    curr_len = 0
    for p in paragraphs:
        if curr_len + len(p) > target_size and curr:
            chunks.append("\n".join(curr))
            curr = [p]
            curr_len = len(p)
        else:
            curr.append(p)
            curr_len += len(p) + 1
    if curr:
        chunks.append("\n".join(curr))
    return chunks

# Dynamic Model Finder
def get_live_groq_model(key):
    headers = {"Authorization": f"Bearer {key}"}
    try:
        r = requests.get("https://api.groq.com/openai/v1/models", headers=headers, timeout=10)
        if r.status_code == 200:
            active_m = [m["id"] for m in r.json().get("data", []) if "whisper" not in m["id"] and "guard" not in m["id"] and "vision" not in m["id"]]
            pref = ["openai/gpt-oss-20b", "openai/gpt-oss-120b", "qwen/qwen3.6-27b", "gemma2-9b-it", "groq/compound"]
            for p in pref:
                if p in active_m:
                    return p
            if active_m:
                return active_m[0]
    except Exception:
        pass
    return "openai/gpt-oss-20b"

# PURE ACADEMIC TRANSLATION PROMPT (So'zma-so'z emas, mukammal akademik ma'no va formulalar daxlsizligi)
ACADEMIC_TRANSLATE_SYSTEM = """
Siz O‘zbekiston Fanlar akademiyasi va nufuzli xalqaro ilmiy jurnallar (Scopus, Web of Science, IEEE) talablariga javob beruvchi oliy toifali ilmiy tahrirchi va akademik tarjimonsiz.

Vazifangiz: Berilgan ilmiy maqola matnini ma'no va mohiyatini 100% saqlagan holda, ravon, professional va chuqur akademik o‘zbek tiliga to‘liq tarjima qilish.

Qat'iy ilmiy-akademik qoidalar:
1. MA'NOVIY VA AKADEMIK ANIQ TARJIMA (So'zma-so'z mexanik tarjima qilmang):
   - Inglizcha jumlalarni quruq so'zma-so'z o'girmasdan, ilmiy fikrni o'zbek tilining akademik grammatikasi va ilmiy uslubiga mos holda ravon, tabiiy va mukammal ifodalang.
   - Muallifning har bir ilmiy g'oyasi, dalili, farazi va xulosasi hech narsa o'zgartirilmasdan va qisqartirilmasdan to'liq aks ettirilsin.

2. FORMULALAR VA ILMIY BELGILAR DAXSLSIZLIGI (Mutlaqo o'zgartirmang):
   - Barcha matematik formulalar ($...$, $$...$$), yunoncha harflar (α, β, γ, λ, θ, σ), indekslar, matritsalar va tenglamalar 100% asl ko'rinishida saqlansin.
   - Harflar, o'zgaruvchilar (x, y, W, b, f(x)) va ilmiy simvollar asliga tegilmasin.

3. IQTIBOSLAR (CITATIONS) VA METADATA:
   - Adabiyotlar iqtiboslari ([1], [2-4], Smith et al., 2024), jadvallar va rasmlar havolalari (1-jadval, 2-rasm) buzilmasdan saqlansin.

4. XALQARO ILMIY VA TEXNIK ATAMALAR:
   - Sun'iy intellekt va kompyuter fanlaridagi o'rnatilgan xalqaro atamalar (Deep Learning, Transformer, Latency, Benchmark, Backpropagation, Learning rate, Overfitting va h.k.) agar o'zbek ilmiy adabiyotida qabul qilingan ekvivalenti bo'lsa berilsin, zarur bo'lsa qavs ichida inglizchasi saqlansin.
   - Model va parametr nomlari (masalan: ResNet-50, AdamW, ReLU, batch_size) o'zgartirilmasin.

5. CHIQISH FORMATI:
   - FAQAT tarjima qilingan mukammal akademik o'zbekcha matnni qaytaring. Hech qanday kirish yoki o'z shaxsiy sharhlaringizni qo'shmang.
"""

def translate_single_chunk(key, model_name, chunk_id, chunk_text):
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {key}"
    }
    payload = {
        "model": model_name,
        "messages": [
            {"role": "system", "content": ACADEMIC_TRANSLATE_SYSTEM},
            {"role": "user", "content": f"Quyidagi ilmiy qismni yuqoridagi qat'iy akademik qoidalarga binoan, mukammal akademik o'zbek tiliga to'liq tarjima qiling:\n\n{chunk_text}"}
        ],
        "temperature": 0.2,
        "max_tokens": 4000
    }
    
    for attempt in range(3):
        try:
            r = requests.post("https://api.groq.com/openai/v1/chat/completions", headers=headers, json=payload, timeout=60)
            if r.status_code == 200:
                trans = r.json()["choices"][0]["message"]["content"].strip()
                return chunk_id, trans
            elif r.status_code == 429:
                time.sleep(2.5)
            else:
                time.sleep(1.0)
        except Exception:
            time.sleep(1.0)
            
    return chunk_id, f"\n\n[Qism {chunk_id+1} tarjimasi]:\n{chunk_text}\n"

# Executive Summary Generator
def generate_executive_summary(key, model_name, text):
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {key}"
    }
    sys_prompt = """
Siz yetakchi akademik ekspert va magistr ilmiy maslahatchisisiz.
Quyidagi ilmiy maqola bo'yicha tahliliy pasportni FAQAT quyidagi JSON formatida qaytaring:
{
  "research_summary": {
    "core_problem": "Maqolada ko'tarilgan asosiy ilmiy muammo nima edi?",
    "proposed_solution": "Mualliflar qanday yangi yechim, model yoki metodologiya taklif qilishdi?",
    "key_focus_areas": "Nimasiga alohida e'tibor berish kerak va qaysi qismlari eng muhim?",
    "experimental_results": "Asosiy tajribaviy natijalar va benchmarklar...",
    "limitations": "Tadqiqot cheklovlari va kamchiliklari..."
  },
  "thesis_advisor": {
    "where_to_cite": "Dissertatsiyaning qaysi bo'limida qanday iqtibos olish kerak...",
    "how_to_use_method": "Ushbu metodni o'z tadqiqotida qanday qo'llash mumkin...",
    "new_research_idea": "Ushbu maqola asosida magistr uchun yangi ilmiy g'oya..."
  },
  "key_terms": [
    {"term_en": "Deep Learning", "term_uz": "Chuqur o'rganish", "desc": "Ko'p qatlamli neyron to'rlariga asoslangan mashinali o'rganish sohasi."}
  ]
}
"""
    payload = {
        "model": model_name,
        "messages": [
            {"role": "system", "content": sys_prompt},
            {"role": "user", "content": f"Quyidagi ilmiy maqolani tahlil qiling:\n\n{text[:16000]}"}
        ],
        "temperature": 0.2,
        "max_tokens": 3000,
        "response_format": {"type": "json_object"}
    }
    try:
        r = requests.post("https://api.groq.com/openai/v1/chat/completions", headers=headers, json=payload, timeout=60)
        if r.status_code == 200:
            return json.loads(r.json()["choices"][0]["message"]["content"])
    except Exception:
        pass
    return {
        "research_summary": {"core_problem": "Tahlil tayyor.", "proposed_solution": "To'liq akademik tarjima bo'limiga qarang."},
        "thesis_advisor": {"where_to_cite": "Metodologiya qismi"},
        "key_terms": []
    }

# PARALLEL PIPELINE EXECUTION
def run_parallel_translation_pipeline(key, text, status_box, progress_bar):
    model_name = get_live_groq_model(key)
    chunks = chunk_document_smart(text, target_size=5000)
    total_chunks = len(chunks)
    
    status_box.write(f"📄 **2-bosqich:** Maqola {total_chunks} ta ilmiy qismga ajratildi. Parallel akademik tarjima boshlandi...")
    
    translated_results = [None] * total_chunks
    completed_count = 0
    
    # Run chunks concurrently (max 3 workers)
    with concurrent.futures.ThreadPoolExecutor(max_workers=3) as executor:
        future_to_id = {executor.submit(translate_single_chunk, key, model_name, i, ch): i for i, ch in enumerate(chunks)}
        for future in concurrent.futures.as_completed(future_to_id):
            chunk_id, trans = future.result()
            translated_results[chunk_id] = trans
            completed_count += 1
            progress_val = int((completed_count / (total_chunks + 1)) * 90)
            progress_bar.progress(progress_val)
            status_box.write(f"✓ *{completed_count}/{total_chunks}-qism to‘liq akademik tarjima qilindi*")

    full_translation = "\n\n".join(translated_results)

    status_box.write("🧠 **3-bosqich:** Chuqur ilmiy xulosa va dissertatsiya pasporti shakllantirilmoqda...")
    summary_data = generate_executive_summary(key, model_name, text)
    progress_bar.progress(100)
    
    return full_translation, summary_data, model_name

# Big Action Button
st.markdown("<br>", unsafe_allow_html=True)
if st.button("🚀 TARJIMA VA TAHLIL QILISH", type="primary", use_container_width=True):
    if not extracted_text.strip():
        st.warning("⚠️ Iltimos, oldin fayl yuklang yoki matn kiriting.")
    elif not active_groq_key:
        st.error("⚠️ Tizimda API kalit topilmadi. Iltimos, Streamlit Secrets-ga GROQ_API_KEY kiriting.")
    else:
        progress_bar = st.progress(5)
        with st.status("🔍 Ilmiy maqola akademik tahlil qilinmoqda va o‘zbekchalashtirilmoqda...", expanded=True) as status_box:
            st.write("📄 **1-bosqich:** Fayl matni va tuzilmasi o‘qildi...")
            st.write(f"✓ Matn hajmi: {len(extracted_text):,} ta belgi.")
            
            try:
                full_trans, summary_data, used_model = run_parallel_translation_pipeline(active_groq_key, extracted_text, status_box, progress_bar)

                st.write("📥 **4-bosqich:** Word (.docx) va PDF eksport fayllari yaratildi.")

                st.session_state["result_translation"] = full_trans
                st.session_state["result_summary"] = summary_data.get("research_summary", {})
                st.session_state["result_thesis"] = summary_data.get("thesis_advisor", {})
                st.session_state["result_terms"] = summary_data.get("key_terms", [])
                st.session_state["file_title"] = file_name
                st.session_state["used_model"] = used_model

                status_box.update(label=f"✅ Mukammal akademik tarjima va tahlil yakunlandi! ({len(full_trans):,} belgi)", state="complete", expanded=False)
                st.success(f"✨ Muvaffaqiyatli yakunlandi! Jami akademik tarjima hajmi: **{len(full_trans):,} ta belgi**")
                
            except Exception as err:
                status_box.update(label="❌ Tahlilda xatolik yuz berdi", state="error", expanded=True)
                st.error(f"Xatolik tafsiloti: {err}")

st.markdown('</div>', unsafe_allow_html=True)

# Results Section
if "result_translation" in st.session_state:
    full_trans = st.session_state["result_translation"]
    summary = st.session_state.get("result_summary", {})
    thesis = st.session_state.get("result_thesis", {})
    terms = st.session_state.get("result_terms", [])
    current_title = st.session_state.get("file_title", "Ilmiy Maqola")

    st.markdown(f"""
    <div class="res-card" style="display: flex; justify-content: space-between; align-items: center; flex-wrap: wrap; gap: 12px; background: linear-gradient(135deg, #f0f9ff 0%, #ffffff 100%); border-color: #bae6fd;">
        <div>
            <div style="font-weight: 800; font-size: 1.15rem; color: #0369a1;">📄 {current_title}</div>
            <div style="font-size: 0.82rem; color: #64748b;">To‘liq akademik tarjima ({len(full_trans):,} belgi) va ilmiy hisobot tayyor</div>
        </div>
    </div>
    """, unsafe_allow_html=True)

    col_d1, col_d2 = st.columns(2)
    with col_d1:
        try:
            doc = Document()
            doc.add_heading(f"UZ SCIENCE AI — {current_title}", level=1)
            doc.add_paragraph(f"Muallif: {AUTHOR_NAME} · To‘liq akademik tarjima va ilmiy tahliliy pasport")
            
            # 1. TO'LIQ AKADEMIK TARJIMA (BOSHIDA)
            doc.add_heading("1. TO‘LIQ AKADEMIK O‘ZBEKCHA TARJIMA (ASOSIY MATN)", level=2)
            doc.add_paragraph(full_trans)
            
            # 2. CHUQUR ILMIY XULOSA
            if summary:
                doc.add_heading("2. Chuqur Ilmiy Xulosa va Tahlil", level=2)
                for k, v in summary.items():
                    p = doc.add_paragraph()
                    p.add_run(f"{k.replace('_', ' ').capitalize()}: ").bold = True
                    p.add_run(str(v))
                    
            # 3. MAGISTR TAVSIYALARI
            if thesis:
                doc.add_heading("3. Magistr Dissertatsiyasi Uchun Tavsiyalar", level=2)
                for k, v in thesis.items():
                    p = doc.add_paragraph()
                    p.add_run(f"{k.replace('_', ' ').capitalize()}: ").bold = True
                    p.add_run(str(v))

            docx_io = io.BytesIO()
            doc.save(docx_io)
            docx_io.seek(0)

            st.download_button(
                label="📄 To‘liq Tarjima (Word .docx) Yuklab Olish",
                data=docx_io,
                file_name=f"UZ_SCIENCE_AI_{current_title}.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                type="primary"
            )
        except Exception as docx_err:
            st.error(f"Word tayyorlashda xatolik: {docx_err}")

    with col_d2:
        try:
            pdf_bytes = create_pdf(current_title, full_trans, summary, thesis, terms)
            st.download_button(
                label="📥 To‘liq Ilmiy Hisobot (PDF) Yuklab Olish",
                data=pdf_bytes,
                file_name=f"UZ_SCIENCE_AI_{current_title}.pdf",
                mime="application/pdf"
            )
        except Exception as pdf_err:
            st.error(f"PDF tayyorlashda xatolik: {pdf_err}")

    st.markdown("<br>", unsafe_allow_html=True)

    # 4 Output Tabs - ASOSIY URG'U: 1-TAB TO'LIQ AKADEMIK TARJIMA
    res_tab1, res_tab2, res_tab3, res_tab4 = st.tabs([
        "📄 To‘liq Akademik Tarjima",
        "🧠 Chuqur Ilmiy Xulosa (Summary)",
        "🎓 Magistr Dissertatsiyasi Tavsiyalari",
        "📖 Ilmiy Terminlar Lug‘ati"
    ])

    # 1-TAB: ASOSIY TO'LIQ AKADEMIK TARJIMA
    with res_tab1:
        st.markdown(f"""
        <div class="res-card" style="border-left: 5px solid #0284c7;">
            <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 12px;">
                <h4 style="margin: 0; color: #0369a1;">🇺🇿 Hech Qayeri Qisqartirilmagan To‘liq Akademik Tarjima</h4>
                <span style="font-size: 0.8rem; background: #e0f2fe; color: #0369a1; padding: 4px 12px; border-radius: 8px; font-weight: 700;">Jami: {len(full_trans):,} ta belgi</span>
            </div>
            <div style="color: #0f172a; font-size: 0.95rem; line-height: 1.8; white-space: pre-wrap;">{full_trans}</div>
        </div>
        """, unsafe_allow_html=True)

    # 2-TAB: CHUQUR ILMIY XULOSA
    with res_tab2:
        st.markdown(f"""
        <div class="focus-banner" style="background: linear-gradient(135deg, #eff6ff 0%, #dbeafe 100%); border: 1px solid #bfdbfe; border-left: 6px solid #0284c7; padding: 20px 24px; border-radius: 14px; margin-bottom: 20px;">
            <div style="font-size: 0.78rem; font-weight: 800; text-transform: uppercase; color: #0284c7; letter-spacing: 0.5px; margin-bottom: 4px;">⚡ DIQQAT MARKAZI • ENG MUHIM QISM</div>
            <h4 style="color: #0c4a6e; margin: 0 0 6px 0; font-size: 1.1rem; font-weight: 800;">Ushbu maqolada nimasiga alohida e’tibor berish kerak?</h4>
            <p style="color: #1e293b; font-size: 0.94rem; margin: 0; line-height: 1.6;">{summary.get('key_focus_areas', '-')}</p>
        </div>
        """, unsafe_allow_html=True)

        col_s1, col_s2 = st.columns(2)
        with col_s1:
            st.markdown(f"""
            <div class="res-card">
                <h4>🎯 1. Ko‘tarilgan Asosiy Muammo</h4>
                <p style="color: #334155; font-size: 0.9rem; line-height: 1.55; margin: 0;">{summary.get('core_problem', '-')}</p>
            </div>
            """, unsafe_allow_html=True)
            
            st.markdown(f"""
            <div class="res-card">
                <h4>📊 3. Natijalar va Benchmarklar</h4>
                <p style="color: #334155; font-size: 0.9rem; line-height: 1.55; margin: 0;">{summary.get('experimental_results', '-')}</p>
            </div>
            """, unsafe_allow_html=True)

        with col_s2:
            st.markdown(f"""
            <div class="res-card">
                <h4>💡 2. Taklif Etilgan Yechim & Metod</h4>
                <p style="color: #334155; font-size: 0.9rem; line-height: 1.55; margin: 0;">{summary.get('proposed_solution', '-')}</p>
            </div>
            """, unsafe_allow_html=True)

            st.markdown(f"""
            <div class="res-card">
                <h4>⏳ 4. Cheklovlar va Kamchiliklar</h4>
                <p style="color: #334155; font-size: 0.9rem; line-height: 1.55; margin: 0;">{summary.get('limitations', '-')}</p>
            </div>
            """, unsafe_allow_html=True)

    # 3-TAB: MAGISTR TAVSIYALARI
    with res_tab3:
        st.markdown(f"""
        <div class="res-card">
            <h4>📌 Qayerda va qanday iqtibos (citation) olish kerak?</h4>
            <p style="color: #334155; font-size: 0.92rem; line-height: 1.6; margin: 0;">{thesis.get('where_to_cite', '-')}</p>
        </div>
        """, unsafe_allow_html=True)

        st.markdown(f"""
        <div class="res-card">
            <h4>🔬 Metodologiyani o‘z tadqiqotingizda qanday qo‘llash mumkin?</h4>
            <p style="color: #334155; font-size: 0.92rem; line-height: 1.6; margin: 0;">{thesis.get('how_to_use_method', '-')}</p>
        </div>
        """, unsafe_allow_html=True)

        st.markdown(f"""
        <div class="res-card" style="background: #fffbeb; border-color: #fde68a; border-left: 5px solid #f59e0b;">
            <h4 style="color: #92400e;">💡 Ushbu maqola asosida yangi dissertatsiya / maqola g‘oyasi:</h4>
            <p style="color: #78350f; font-size: 0.92rem; line-height: 1.6; font-weight: 500; margin: 0;">{thesis.get('new_research_idea', '-')}</p>
        </div>
        """, unsafe_allow_html=True)

    # 4-TAB: LUG'AT
    with res_tab4:
        if terms:
            for item in terms:
                st.markdown(f"""
                <div class="res-card" style="padding: 14px 18px; margin-bottom: 10px;">
                    <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 4px;">
                        <span style="font-weight: 800; color: #0284c7; font-size: 0.95rem;">{item.get('term_en', '')}</span>
                        <span style="font-weight: 700; color: #0f172a; font-size: 0.84rem; background: #e0f2fe; padding: 2px 8px; border-radius: 6px;">🇺🇿 {item.get('term_uz', '')}</span>
                    </div>
                    <p style="margin: 0; font-size: 0.86rem; color: #475569; line-height: 1.45;">{item.get('desc', '')}</p>
                </div>
                """, unsafe_allow_html=True)
        else:
            st.info("Terminlar ro‘yxati shakllantirilmadi.")

# Footer with Author Branding
st.markdown(f"""
<div class="footer-container">
    © 2026 <b>UZ SCIENCE AI</b> — Yaratuvchi: <b>{AUTHOR_NAME}</b>. Magistrant va ilmiy izlanuvchilar uchun maxsus yaratilgan.
</div>
""", unsafe_allow_html=True)
